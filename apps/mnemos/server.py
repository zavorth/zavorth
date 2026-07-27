"""
Zavorth local memory engine.

Exposes semantic search, metadata scanning, and file indexing tools over MCP.

Architecture:
  - ChromaDB: local vector database.
  - SentenceTransformers: lightweight local model for vector generation.
  - Watchdog: asynchronous sentinel for automatic vault indexing.

Security:
  - This server runs in an isolated Docker container.
  - Scan directories are mounted as READ-ONLY volumes.
  - No host file is accessed without explicit user authorization.
"""

import os
import sys
import json
import logging
import hashlib
import mimetypes
import re
import zipfile
from html import unescape
from pathlib import Path
from typing import Any, Optional

from mcp.server import Server
from mcp.server.stdio import stdio_server
from mcp.types import TextContent, Tool

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

CHROMA_DB_DIR = os.environ.get("MNEMOS_CHROMA_DB_DIR", "/app/data/vector_db")
SCAN_VOLUMES_BASE = os.environ.get("MNEMOS_SCAN_VOLUMES", "/scan_volumes")
VAULT_DIR = os.environ.get("MNEMOS_VAULT_DIR", "/app/data/vault")
EMBEDDING_MODEL = os.environ.get("MNEMOS_EMBEDDING_MODEL", "all-MiniLM-L6-v2")
MAX_SCAN_DEPTH = int(os.environ.get("MNEMOS_MAX_SCAN_DEPTH", "5"))
MAX_RESULTS = int(os.environ.get("MNEMOS_MAX_RESULTS", "10"))

logging.basicConfig(
    level=logging.INFO,
    format="[Mnemos] %(asctime)s %(levelname)s %(message)s",
    stream=sys.stderr,  # stdout is reserved for MCP
)
logger = logging.getLogger("mnemos")

# ---------------------------------------------------------------------------
# Motor Vetorial (Lazy-loaded para saving de RAM)
# ---------------------------------------------------------------------------

_chroma_client = None
_chroma_collection = None
_embedding_function = None


def _get_embedding_function():
    """Carrega o modelo de embedding local sob demanda."""
    global _embedding_function
    if _embedding_function is None:
        try:
            from chromadb.utils.embedding_functions import SentenceTransformerEmbeddingFunction
            logger.info(f"Carregando modelo de embedding: {EMBEDDING_MODEL}")
            _embedding_function = SentenceTransformerEmbeddingFunction(
                model_name=EMBEDDING_MODEL,
            )
            logger.info("Modelo de embedding loaded com success.")
        except Exception as e:
            logger.error(f"Failed to load embedding model: {e}")
            raise
    return _embedding_function


def _get_collection():
    """Inicializa o ChromaDB e retorna a collection principal."""
    global _chroma_client, _chroma_collection
    if _chroma_collection is None:
        import chromadb

        persist_dir = Path(CHROMA_DB_DIR)
        persist_dir.mkdir(parents=True, exist_ok=True)

        logger.info(f"Inicializando ChromaDB em: {persist_dir}")
        _chroma_client = chromadb.PersistentClient(path=str(persist_dir))
        _chroma_collection = _chroma_client.get_or_create_collection(
            name="mnemos_vault",
            embedding_function=_get_embedding_function(),
            metadata={"hnsw:space": "cosine"},
        )
        logger.info(
            f"Collection 'mnemos_vault' ready. Documents current: {_chroma_collection.count()}"
        )
    return _chroma_collection


# ---------------------------------------------------------------------------
# Text extraction
# ---------------------------------------------------------------------------

TEXT_EXTENSIONS = {
    ".txt", ".md", ".py", ".ts", ".tsx", ".js", ".jsx", ".json", ".csv",
    ".log", ".html", ".htm", ".xml", ".yaml", ".yml", ".toml", ".ini",
    ".cfg", ".tex", ".rst", ".sql", ".java", ".cs", ".go", ".rs", ".php",
}
PDF_EXTENSIONS = {".pdf"}
DOCX_EXTENSIONS = {".docx"}
PPTX_EXTENSIONS = {".pptx"}
XLSX_EXTENSIONS = {".xlsx", ".xlsm"}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}
AUDIO_VIDEO_EXTENSIONS = {".mp3", ".wav", ".m4a", ".ogg", ".mp4", ".mov", ".mkv", ".webm"}

MAX_UNDERSTANDING_PREVIEW_CHARS = int(os.environ.get("MNEMOS_MAX_UNDERSTANDING_PREVIEW_CHARS", "4000"))
SUPPORTED_EXTENSIONS = TEXT_EXTENSIONS


def _safe_resolve_path(value: str) -> Path:
    try:
        return Path(value).resolve(strict=False)
    except Exception:
        return Path(value).absolute()


def _authorized_roots() -> list[Path]:
    roots: list[Path] = []
    for raw in [SCAN_VOLUMES_BASE, VAULT_DIR]:
        if not raw:
            continue
        root = _safe_resolve_path(raw)
        if root.exists():
            roots.append(root)
    return roots


def _is_relative_to(candidate: Path, root: Path) -> bool:
    try:
        candidate.relative_to(root)
        return True
    except ValueError:
        return False


def _resolve_authorized_file(file_path: str) -> tuple[Optional[Path], Optional[dict[str, Any]]]:
    raw = str(file_path or "").strip()
    if not raw:
        return None, {"status": "error", "error": "File path was not specified."}
    fp = _safe_resolve_path(raw)
    roots = _authorized_roots()
    if not roots or not any(_is_relative_to(fp, root) for root in roots):
        return None, {
            "status": "blocked",
            "error": "File is outside authorized Mnemos volumes.",
            "file_name": Path(raw).name,
            "allowed_roots": [str(root) for root in roots],
            "indexable_text": "",
        }
    if not fp.exists() or not fp.is_file():
        return None, {
            "status": "error",
            "error": f"File not found: {Path(raw).name}",
            "file_name": Path(raw).name,
            "indexable_text": "",
        }
    return fp, None


def _extract_text(file_path: str) -> Optional[str]:
    """Extracts text from a file, supporting plain text and PDFs."""
    fp, error = _resolve_authorized_file(file_path)
    if error or fp is None:
        return None

    ext = fp.suffix.lower()

    if ext in SUPPORTED_EXTENSIONS:
        try:
            return fp.read_text(encoding="utf-8", errors="replace")
        except Exception as e:
            logger.warning(f"Failed to read text from {file_path}: {e}")
            return None

    if ext in PDF_EXTENSIONS:
        try:
            import subprocess
            # Use pdftotext when available in the container.
            result = subprocess.run(
                ["pdftotext", file_path, "-"],
                capture_output=True, text=True, timeout=30,
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        except FileNotFoundError:
            pass
        except Exception as e:
            logger.warning(f"pdftotext failed for {file_path}: {e}")

        # Fallback: pypdf
        try:
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            pages = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    pages.append(text)
            return "\n\n".join(pages) if pages else None
        except ImportError:
            logger.warning("pypdf is not installed; PDF text extraction is unavailable.")
            return None
        except Exception as e:
            logger.warning(f"pypdf failed for {file_path}: {e}")
            return None

    return None


def _clean_text(text: str) -> str:
    return re.sub(r"\n{3,}", "\n\n", unescape(str(text or "").replace("\x00", " "))).strip()


def _strip_xml_text(raw: str) -> str:
    return _clean_text(re.sub(r"<[^>]+>", " ", raw))


def _append_text_part(parts: list[dict[str, Any]], kind: str, text: Optional[str], **metadata: Any) -> None:
    cleaned = _clean_text(text or "")
    if cleaned:
        parts.append({"kind": kind, "text": cleaned, **metadata})


def _classify_document_type(ext: str) -> str:
    if ext in TEXT_EXTENSIONS:
        return "text"
    if ext in PDF_EXTENSIONS:
        return "pdf"
    if ext in DOCX_EXTENSIONS:
        return "docx"
    if ext in PPTX_EXTENSIONS:
        return "pptx"
    if ext in XLSX_EXTENSIONS:
        return "xlsx"
    if ext in IMAGE_EXTENSIONS:
        return "image"
    if ext in AUDIO_VIDEO_EXTENSIONS:
        return "audio-video"
    return "unknown"


def _extract_plain_text(fp: Path, parts: list[dict[str, Any]], extractors: list[str], warnings: list[str]) -> None:
    try:
        _append_text_part(parts, "text", fp.read_text(encoding="utf-8", errors="replace"))
        extractors.append("plain-text")
    except Exception as e:
        warnings.append(f"plain-text failed: {e}")


def _extract_pdf_text_universal(fp: Path, parts: list[dict[str, Any]], extractors: list[str], warnings: list[str]) -> bool:
    before = len(parts)
    try:
        import subprocess
        result = subprocess.run(
            ["pdftotext", str(fp), "-"],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            _append_text_part(parts, "pdf-text", result.stdout)
            extractors.append("pdf-text:pdftotext")
    except FileNotFoundError:
        warnings.append("pdftotext not installed")
    except Exception as e:
        warnings.append(f"pdftotext failed: {e}")

    if len(parts) == before:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(fp))
            pages = []
            for page_index, page in enumerate(reader.pages):
                text = page.extract_text()
                if text:
                    pages.append(f"[page {page_index + 1}]\n{text}")
            if pages:
                _append_text_part(parts, "pdf-text", "\n\n".join(pages))
                extractors.append("pdf-text:pypdf")
        except ImportError:
            warnings.append("pypdf not installed")
        except Exception as e:
            warnings.append(f"pypdf failed: {e}")

    return len(parts) > before


def _extract_pdf_ocr(fp: Path, parts: list[dict[str, Any]], extractors: list[str], warnings: list[str]) -> bool:
    before = len(parts)
    try:
        from pdf2image import convert_from_path
        import pytesseract

        pages = convert_from_path(str(fp), dpi=200, first_page=1, last_page=20)
        page_texts = []
        for index, image in enumerate(pages):
            text = pytesseract.image_to_string(image, lang=os.environ.get("MNEMOS_OCR_LANG", "por+eng"))
            if text and text.strip():
                page_texts.append(f"[ocr page {index + 1}]\n{text.strip()}")
        if page_texts:
            _append_text_part(parts, "pdf-ocr", "\n\n".join(page_texts))
            extractors.append("pdf-ocr:tesseract")
    except ImportError:
        warnings.append("pdf2image/pytesseract not installed")
    except Exception as e:
        warnings.append(f"pdf OCR failed: {e}")
    return len(parts) > before


def _extract_docx(fp: Path, parts: list[dict[str, Any]], extractors: list[str], warnings: list[str], tables: list[dict[str, Any]]) -> None:
    try:
        from docx import Document
        document = Document(str(fp))
        paragraphs = [p.text for p in document.paragraphs if p.text and p.text.strip()]
        _append_text_part(parts, "docx-paragraphs", "\n\n".join(paragraphs))
        for table_index, table in enumerate(document.tables):
            rows = []
            for row in table.rows:
                cells = [_clean_text(cell.text) for cell in row.cells]
                if any(cells):
                    rows.append("\t".join(cells))
            if rows:
                table_text = "\n".join(rows)
                tables.append({"index": table_index, "rows": len(rows), "preview": table_text[:1000]})
                _append_text_part(parts, "docx-table", table_text, table_index=table_index)
        extractors.append("docx:python-docx")
        return
    except ImportError:
        warnings.append("python-docx not installed")
    except Exception as e:
        warnings.append(f"python-docx failed: {e}")

    try:
        with zipfile.ZipFile(fp) as archive:
            xml = archive.read("word/document.xml").decode("utf-8", errors="replace")
            _append_text_part(parts, "docx-zip-fallback", _strip_xml_text(xml))
            extractors.append("docx:zip-fallback")
    except Exception as e:
        warnings.append(f"docx zip fallback failed: {e}")


def _extract_xlsx(fp: Path, parts: list[dict[str, Any]], extractors: list[str], warnings: list[str], tables: list[dict[str, Any]]) -> None:
    try:
        from openpyxl import load_workbook
        workbook = load_workbook(str(fp), data_only=True, read_only=True)
        for sheet in workbook.worksheets:
            rows = []
            for row_index, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if row_index > 300:
                    warnings.append(f"sheet {sheet.title} truncated at 300 rows")
                    break
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    rows.append("\t".join(values))
            if rows:
                table_text = "\n".join(rows)
                tables.append({"sheet": sheet.title, "rows": len(rows), "preview": table_text[:1000]})
                _append_text_part(parts, "xlsx-sheet", f"[sheet {sheet.title}]\n{table_text}", sheet=sheet.title)
        extractors.append("xlsx:openpyxl")
    except ImportError:
        warnings.append("openpyxl not installed")
    except Exception as e:
        warnings.append(f"xlsx extraction failed: {e}")


def _extract_pptx(fp: Path, parts: list[dict[str, Any]], extractors: list[str], warnings: list[str], tables: list[dict[str, Any]]) -> None:
    try:
        from pptx import Presentation
        presentation = Presentation(str(fp))
        for slide_index, slide in enumerate(presentation.slides, start=1):
            slide_chunks = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_chunks.append(shape.text)
                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        cells = [_clean_text(cell.text) for cell in row.cells]
                        if any(cells):
                            rows.append("\t".join(cells))
                    if rows:
                        table_text = "\n".join(rows)
                        tables.append({"slide": slide_index, "rows": len(rows), "preview": table_text[:1000]})
                        slide_chunks.append(table_text)
            _append_text_part(parts, "pptx-slide", "\n\n".join(slide_chunks), slide=slide_index)
        extractors.append("pptx:python-pptx")
    except ImportError:
        warnings.append("python-pptx not installed")
    except Exception as e:
        warnings.append(f"pptx extraction failed: {e}")


def _extract_image_ocr(fp: Path, parts: list[dict[str, Any]], extractors: list[str], warnings: list[str], visual: dict[str, Any]) -> bool:
    before = len(parts)
    try:
        from PIL import Image
        import pytesseract
        with Image.open(fp) as image:
            visual["dimensions"] = {"width": image.width, "height": image.height}
            visual["mode"] = image.mode
            text = pytesseract.image_to_string(image, lang=os.environ.get("MNEMOS_OCR_LANG", "por+eng"))
            _append_text_part(parts, "image-ocr", text)
        extractors.append("image-ocr:tesseract")
    except ImportError:
        warnings.append("Pillow/pytesseract not installed")
    except Exception as e:
        warnings.append(f"image OCR failed: {e}")
    return len(parts) > before


def _build_universal_understanding(file_path: str) -> dict[str, Any]:
    fp, error = _resolve_authorized_file(file_path)
    if error or fp is None:
        return error or {"status": "error", "error": "File is not authorized.", "indexable_text": ""}

    ext = fp.suffix.lower()
    document_type = _classify_document_type(ext)
    parts: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []
    warnings: list[str] = []
    extractors: list[str] = []
    visual: dict[str, Any] = {}
    needs_vision_provider = False
    needs_transcription = False

    if document_type == "text":
        _extract_plain_text(fp, parts, extractors, warnings)
    elif document_type == "pdf":
        has_text = _extract_pdf_text_universal(fp, parts, extractors, warnings)
        if not has_text:
            _extract_pdf_ocr(fp, parts, extractors, warnings)
        needs_vision_provider = len(parts) == 0
    elif document_type == "docx":
        _extract_docx(fp, parts, extractors, warnings, tables)
    elif document_type == "xlsx":
        _extract_xlsx(fp, parts, extractors, warnings, tables)
        needs_vision_provider = True
        warnings.append("Charts/images in spreadsheets require a governed multimodal provider pass.")
    elif document_type == "pptx":
        _extract_pptx(fp, parts, extractors, warnings, tables)
        needs_vision_provider = True
        warnings.append("Visual slide layout, diagrams and embedded images require a governed multimodal provider pass.")
    elif document_type == "image":
        _extract_image_ocr(fp, parts, extractors, warnings, visual)
        needs_vision_provider = True
        warnings.append("Image semantics, charts and diagrams require a governed multimodal provider pass.")
    elif document_type == "audio-video":
        needs_transcription = True
        warnings.append("Audio/video transcription is not local in Mnemos yet; route through a governed transcription capability.")
    else:
        warnings.append(f"Unsupported extension: {ext or 'none'}")

    indexable_text = _clean_text("\n\n".join(
        f"[{part.get('kind')}]\n{part.get('text')}" for part in parts if part.get("text")
    ))
    if not indexable_text and document_type in {"image", "audio-video"}:
        indexable_text = _clean_text(
            f"{document_type} file {fp.name}. No local text was extracted. "
            "A governed multimodal/transcription provider pass is required for semantic understanding."
        )

    confidence = 0.0
    if indexable_text and any("ocr" in extractor for extractor in extractors):
        confidence = 0.72
    elif indexable_text:
        confidence = 0.9
    elif needs_vision_provider or needs_transcription:
        confidence = 0.35

    file_stat = fp.stat()
    receipt_id = hashlib.sha256(f"{file_path}:{file_stat.st_mtime_ns}:{file_stat.st_size}".encode("utf-8")).hexdigest()[:16]
    chunks_estimate = len(_chunk_text(indexable_text)) if indexable_text else 0
    status = "success" if indexable_text else "partial" if (needs_vision_provider or needs_transcription) else "unsupported"

    return {
        "status": status,
        "file_name": fp.name,
        "file_path": str(fp),
        "extension": ext,
        "mime_type": mimetypes.guess_type(fp.name)[0] or "application/octet-stream",
        "size_bytes": file_stat.st_size,
        "document_type": document_type,
        "extractors": extractors,
        "warnings": warnings,
        "tables": tables,
        "visual": visual,
        "capabilities": {
            "text_extracted": bool(indexable_text),
            "tables_extracted": len(tables) > 0,
            "ocr_used": any("ocr" in extractor for extractor in extractors),
            "vision_required": needs_vision_provider,
            "transcription_required": needs_transcription,
            "external_provider_used": False,
        },
        "confidence": confidence,
        "chunks_estimate": chunks_estimate,
        "text_preview": indexable_text[:MAX_UNDERSTANDING_PREVIEW_CHARS],
        "extracted_text_length": len(indexable_text),
        "indexable_text": indexable_text,
        "receipt": {
            "id": f"mnemos-understanding:{receipt_id}",
            "kind": "mnemos-universal-file-understanding",
            "file": fp.name,
            "document_type": document_type,
            "extractors": extractors,
            "chunks_estimate": chunks_estimate,
            "tables": len(tables),
            "vision_required": needs_vision_provider,
            "transcription_required": needs_transcription,
            "external_provider_used": False,
            "confidence": confidence,
        },
    }


def _chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> list[str]:
    """Splits text into overlapping chunks for better search recall."""
    if len(text) <= chunk_size:
        return [text]

    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk.strip())
        start = end - overlap

    return chunks


# ---------------------------------------------------------------------------
# server MCP
# ---------------------------------------------------------------------------

server = Server("mnemos-cognitive-engine")


@server.list_tools()
async def list_tools() -> list[Tool]:
    """Declares the tools Mnemos exposes to Zavorth."""
    return [
        Tool(
            name="search_memory",
            description=(
                "Searches the local Mnemos vector memory database for "
                "knowledge fragments semantically relevant to the query. "
                "Use it to find notes, PDF excerpts, and content previously indexed by the user."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "query": {
                        "type": "string",
                        "description": "The semantic search question or phrase.",
                    },
                    "limit": {
                        "type": "integer",
                        "description": "Maximum number of results.",
                        "default": 5,
                    },
                },
                "required": ["query"],
            },
        ),
        Tool(
            name="scan_local_metadata",
            description=(
                "Scans user-authorized directories for files whose names contain "
                "the provided keywords. Lightweight metadata search without reading file contents."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "keywords": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "List of keywords for filtering file names.",
                    },
                    "extensions": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional file extensions to filter, for example ['.pdf', '.docx'].",
                    },
                },
                "required": ["keywords"],
            },
        ),
        Tool(
            name="index_file",
            description=(
                "Reads, extracts text, and indexes a specific file in the vector database. "
                "Use after the user confirms they want to index a file found by metadata scan."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Absolute path of the file to index inside mounted volumes.",
                    },
                    "tags": {
                        "type": "array",
                        "items": {"type": "string"},
                        "description": "Optional tags for categorizing the document.",
                    },
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="understand_file",
            description=(
                "Analyzes an authorized file with the Mnemos Universal File Understanding pipeline without indexing: text, PDF, OCR, DOCX, XLSX, PPTX, image, tables, and limit receipt."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "file_path": {
                        "type": "string",
                        "description": "Absolute file path inside authorized volumes.",
                    },
                    "include_text": {
                        "type": "boolean",
                        "description": "If true, includes the full extracted text. Default false.",
                        "default": False,
                    },
                },
                "required": ["file_path"],
            },
        ),
        Tool(
            name="vault_status",
            description=(
                "Returns memory vault statistics: indexed document count, disk space, monitored directories, and related status."
            ),
            inputSchema={
                "type": "object",
                "properties": {},
            },
        ),
        Tool(
            name="delete_memory",
            description=(
                "Removes a specific document from the vector database by ID or source file name."
            ),
            inputSchema={
                "type": "object",
                "properties": {
                    "document_id": {
                        "type": "string",
                        "description": "Document ID in the vector database.",
                    },
                    "source_file": {
                        "type": "string",
                        "description": "Source file name used to remove all associated chunks.",
                    },
                },
            },
        ),
    ]


@server.call_tool()
async def call_tool(name: str, arguments: dict) -> list[TextContent]:
    """Routes the call to the correct tool implementation."""
    try:
        if name == "search_memory":
            return await _handle_search_memory(arguments)
        elif name == "scan_local_metadata":
            return await _handle_scan_local_metadata(arguments)
        elif name == "index_file":
            return await _handle_index_file(arguments)
        elif name == "understand_file":
            return await _handle_understand_file(arguments)
        elif name == "vault_status":
            return await _handle_vault_status(arguments)
        elif name == "delete_memory":
            return await _handle_delete_memory(arguments)
        else:
            return [TextContent(type="text", text=f"Unknown tool: {name}")]
    except Exception as e:
        logger.error(f"Error executing tool '{name}': {e}", exc_info=True)
        return [TextContent(type="text", text=f"Internal Mnemos error: {str(e)}")]


# ---------------------------------------------------------------------------
# Tool implementations
# ---------------------------------------------------------------------------


async def _handle_search_memory(args: dict) -> list[TextContent]:
    """Semantic search in the vector database."""
    query = str(args.get("query", "")).strip()
    limit = int(args.get("limit", 5))

    if not query:
        return [TextContent(type="text", text=json.dumps({"error": "Empty query."}))]

    collection = _get_collection()
    if collection.count() == 0:
        return [TextContent(type="text", text=json.dumps({
            "hits": [],
            "total_documents": 0,
            "message": "The vault is empty. No document has been indexed yet.",
        }))]

    results = collection.query(
        query_texts=[query],
        n_results=min(limit, MAX_RESULTS),
        include=["documents", "metadatas", "distances"],
    )

    hits = []
    if results and results.get("documents"):
        for i, doc in enumerate(results["documents"][0]):
            metadata = results["metadatas"][0][i] if results.get("metadatas") else {}
            distance = results["distances"][0][i] if results.get("distances") else None
            hits.append({
                "text": doc,
                "source": metadata.get("source", "unknown"),
                "chunk_index": metadata.get("chunk_index", 0),
                "relevance": round(1.0 - (distance or 0.0), 4),
            })

    response = {
        "hits": hits,
        "total_documents": collection.count(),
        "message": f"Found {len(hits)} relevant fragment(s)." if hits else "No result found in the vault.",
    }

    return [TextContent(type="text", text=json.dumps(response, ensure_ascii=False))]


async def _handle_scan_local_metadata(args: dict) -> list[TextContent]:
    """Lightweight metadata scan in authorized directories."""
    keywords = args.get("keywords", [])
    extensions = args.get("extensions", [])

    if not keywords:
        return [TextContent(type="text", text=json.dumps({"error": "No keyword provided."}))]

    # Normalize keywords and extensions.
    keywords_lower = [kw.lower() for kw in keywords]
    extensions_lower = [ext.lower() if ext.startswith(".") else f".{ext.lower()}" for ext in extensions]

    found_files = []
    scan_base = Path(SCAN_VOLUMES_BASE)

    if not scan_base.exists():
        return [TextContent(type="text", text=json.dumps({
            "results": [],
            "scanned_path": str(scan_base),
            "message": "No authorized scan directory is mounted.",
        }))]

    # Scan with bounded depth.
    def scan_dir(directory: Path, depth: int = 0):
        if depth > MAX_SCAN_DEPTH:
            return
        try:
            for entry in directory.iterdir():
                if entry.is_file():
                    name_lower = entry.name.lower()
                    # Filter by extension when specified.
                    if extensions_lower and entry.suffix.lower() not in extensions_lower:
                        continue
                    # Match by keyword in the file name.
                    if any(kw in name_lower for kw in keywords_lower):
                        found_files.append({
                            "name": entry.name,
                            "path": str(entry),
                            "size_bytes": entry.stat().st_size,
                            "extension": entry.suffix,
                        })
                elif entry.is_dir() and not entry.name.startswith("."):
                    scan_dir(entry, depth + 1)
        except PermissionError:
            pass  # Silently ignore directories without permission.

    scan_dir(scan_base)

    response = {
        "results": found_files[:MAX_RESULTS],
        "total_found": len(found_files),
        "scanned_path": str(scan_base),
        "message": f"Found {len(found_files)} matching file(s)." if found_files else "No file found in authorized directories.",
    }

    return [TextContent(type="text", text=json.dumps(response, ensure_ascii=False))]


async def _handle_understand_file(args: dict) -> list[TextContent]:
    """Preview Universal File Understanding without indexing."""
    file_path = str(args.get("file_path", "")).strip()
    include_text = bool(args.get("include_text", False))

    understanding = _build_universal_understanding(file_path)
    response = dict(understanding)
    if not include_text:
        response.pop("indexable_text", None)

    return [TextContent(type="text", text=json.dumps(response, ensure_ascii=False))]


async def _handle_index_file(args: dict) -> list[TextContent]:
    """File indexing in the vector database."""
    file_path = str(args.get("file_path", "")).strip()
    tags = args.get("tags", [])

    fp, path_error = _resolve_authorized_file(file_path)
    if path_error or fp is None:
        return [TextContent(type="text", text=json.dumps(path_error, ensure_ascii=False))]

    # Text extraction
    understanding = _build_universal_understanding(str(fp))
    if understanding.get("status") == "error":
        return [TextContent(type="text", text=json.dumps({
            "error": understanding.get("error"),
            "file_path": file_path,
        }, ensure_ascii=False))]

    text = str(understanding.get("indexable_text") or "")
    if not text or not text.strip():
        return [TextContent(type="text", text=json.dumps({
            "error": f"Could not extract text from file: {fp.name}",
            "extension": fp.suffix,
            "understanding_receipt": understanding.get("receipt"),
            "warnings": understanding.get("warnings", []),
        }))]

    # Split into chunks
    chunks = _chunk_text(text)
    if not chunks:
        return [TextContent(type="text", text=json.dumps({"error": "No text chunk was generated."}))]

    # Index in ChromaDB
    collection = _get_collection()

    source_hash = hashlib.sha256(file_path.encode("utf-8")).hexdigest()[:10]
    ids = [f"{fp.stem}__{source_hash}__chunk_{i}" for i in range(len(chunks))]
    metadatas = [{
        "source": fp.name,
        "source_path": file_path,
        "chunk_index": i,
        "total_chunks": len(chunks),
        "tags": ",".join(tags) if tags else "",
        "extension": fp.suffix,
        "document_type": understanding.get("document_type", "unknown"),
        "extractors": ",".join(understanding.get("extractors", [])),
        "understanding_receipt_id": (understanding.get("receipt") or {}).get("id", ""),
        "vision_required": str((understanding.get("capabilities") or {}).get("vision_required", False)).lower(),
        "transcription_required": str((understanding.get("capabilities") or {}).get("transcription_required", False)).lower(),
        "confidence": str(understanding.get("confidence", 0.0)),
    } for i in range(len(chunks))]

    # Upsert to allow re-indexing
    collection.upsert(
        ids=ids,
        documents=chunks,
        metadatas=metadatas,
    )

    logger.info(f"Indexed file: {fp.name} ({len(chunks)} chunks)")

    response = {
        "status": "success",
        "file_name": fp.name,
        "chunks_indexed": len(chunks),
        "total_vault_documents": collection.count(),
        "understanding_receipt": understanding.get("receipt"),
        "document_type": understanding.get("document_type"),
        "extractors": understanding.get("extractors", []),
        "capabilities": understanding.get("capabilities", {}),
        "warnings": understanding.get("warnings", []),
        "message": f"'{fp.name}' indexed successfully ({len(chunks)} fragments).",
    }

    return [TextContent(type="text", text=json.dumps(response, ensure_ascii=False))]


async def _handle_vault_status(args: dict) -> list[TextContent]:
    """Returns memory vault statistics."""
    collection = _get_collection()

    # Calculate DB disk usage
    db_path = Path(CHROMA_DB_DIR)
    db_size = sum(f.stat().st_size for f in db_path.rglob("*") if f.is_file()) if db_path.exists() else 0

    # Check mounted volumes
    scan_base = Path(SCAN_VOLUMES_BASE)
    mounted_volumes = []
    if scan_base.exists():
        mounted_volumes = [d.name for d in scan_base.iterdir() if d.is_dir()]

    response = {
        "total_documents": collection.count(),
        "db_size_mb": round(db_size / (1024 * 1024), 2),
        "db_path": str(db_path),
        "embedding_model": EMBEDDING_MODEL,
        "mounted_scan_volumes": mounted_volumes,
        "vault_dir": VAULT_DIR,
        "max_scan_depth": MAX_SCAN_DEPTH,
        "universal_file_understanding": {
            "enabled": True,
            "local_extractors": [
                "plain-text",
                "pdf-text",
                "pdf-ocr",
                "docx",
                "xlsx",
                "pptx",
                "image-ocr",
            ],
            "external_provider_used": False,
        },
    }

    return [TextContent(type="text", text=json.dumps(response, ensure_ascii=False))]


async def _handle_delete_memory(args: dict) -> list[TextContent]:
    """Removes documents from the vector database."""
    document_id = str(args.get("document_id", "")).strip()
    source_file = str(args.get("source_file", "")).strip()

    if not document_id and not source_file:
        return [TextContent(type="text", text=json.dumps({"error": "Specify document_id or source_file."}))]

    collection = _get_collection()

    if document_id:
        try:
            collection.delete(ids=[document_id])
            return [TextContent(type="text", text=json.dumps({
                "status": "success",
                "deleted_id": document_id,
                "message": f"Document '{document_id}' removed from the vault.",
            }))]
        except Exception as e:
            return [TextContent(type="text", text=json.dumps({"error": str(e)}))]

    if source_file:
        # Find all chunks for this source file
        results = collection.get(
            where={"source": source_file},
            include=["metadatas"],
        )
        if results and results.get("ids"):
            collection.delete(ids=results["ids"])
            return [TextContent(type="text", text=json.dumps({
                "status": "success",
                "source_file": source_file,
                "chunks_deleted": len(results["ids"]),
                "message": f"All {len(results['ids'])} chunks from '{source_file}' were removed.",
            }))]
        else:
            return [TextContent(type="text", text=json.dumps({
                "status": "not_found",
                "source_file": source_file,
                "message": f"No chunk found for file '{source_file}'.",
            }))]


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

async def main():
    """Initializes the MCP server through stdin/stdout."""
    logger.info("Starting Mnemos Cognitive Engine (MCP Server)...")
    logger.info(f"ChromaDB dir: {CHROMA_DB_DIR}")
    logger.info(f"Scan volumes: {SCAN_VOLUMES_BASE}")
    logger.info(f"Embedding model: {EMBEDDING_MODEL}")

    async with stdio_server() as (read_stream, write_stream):
        await server.run(
            read_stream,
            write_stream,
            server.create_initialization_options(),
        )


if __name__ == "__main__":
    import asyncio
    asyncio.run(main())
